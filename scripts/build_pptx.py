#!/usr/bin/env python3
"""Build the AltaJuris PowerPoint presentation from the manual.

The script reads `docs/manual/*.md` and builds a deck covering all current
features. Adding content to the manual automatically flows into the pptx
on the next build.

Usage:
    python3 scripts/build_pptx.py
    # output: build/AltaJuris-Apresentacao.pptx
"""

from __future__ import annotations

import datetime as _dt
import pathlib
import re
from dataclasses import dataclass, field

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = pathlib.Path(__file__).resolve().parent.parent
MANUAL_DIR = ROOT / "docs" / "manual"
OUT_DIR = ROOT / "build"
STATIC_DIR = ROOT / "static" / "downloads"
OUT_FILE = OUT_DIR / "AltaJuris-Apresentacao.pptx"

# Brand palette
NAVY = RGBColor(0x1E, 0x3A, 0x5F)
GOLD = RGBColor(0xD4, 0xA5, 0x37)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
SOFT_BG = RGBColor(0xF5, 0xF5, 0xF7)
INDIGO = RGBColor(0x63, 0x66, 0xF1)
EMERALD = RGBColor(0x10, 0xB9, 0x81)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)


# ─── Manual parsing ────────────────────────────────────────────────────────

@dataclass
class ManualSection:
    title: str
    description: str
    bullets: list[str] = field(default_factory=list)


def parse_manual(path: pathlib.Path) -> ManualSection:
    """Extract title/description/main bullets from a manual markdown file."""
    raw = path.read_text(encoding="utf-8")

    # Frontmatter
    title = path.stem.replace("-", " ").title()
    description = ""
    fm = re.match(r"---\n(.*?)\n---\n", raw, re.S)
    body = raw
    if fm:
        block = fm.group(1)
        m = re.search(r"^title:\s*(.+)$", block, re.M)
        if m:
            title = m.group(1).strip().strip("\"'")
        m = re.search(r"^description:\s*(.+)$", block, re.M)
        if m:
            description = m.group(1).strip().strip("\"'")
        body = raw[fm.end() :]

    # Collect bullets from the first few sections
    bullets: list[str] = []
    for line in body.splitlines():
        ln = line.strip()
        if ln.startswith("- ") and len(bullets) < 6:
            text = ln[2:]
            text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)  # drop bold markers
            text = re.sub(r"`(.+?)`", r"\1", text)  # drop inline code
            text = re.sub(r"\[(.+?)\]\((.+?)\)", r"\1", text)  # drop links
            text = re.sub(r"<[^>]+>", "", text)
            if text:
                bullets.append(text)
    return ManualSection(title=title.strip(), description=description.strip(), bullets=bullets)


# ─── Slide helpers ─────────────────────────────────────────────────────────

def _fill_bg(slide, color=NAVY):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color


def _add_rect(slide, left, top, width, height, color, line=False):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if not line:
        shp.line.fill.background()
    return shp


def _add_text(slide, left, top, width, height, text, *, size=18, color=WHITE, bold=False,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    lines = text.split("\n") if "\n" in text else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = "Calibri"
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
    return tb


def _add_bullets(slide, left, top, width, height, items, *, size=16, color=WHITE,
                 bullet_color=GOLD):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        r = p.add_run()
        r.text = "▸ "
        r.font.name = "Calibri"
        r.font.size = Pt(size)
        r.font.color.rgb = bullet_color
        r.font.bold = True
        r = p.add_run()
        r.text = item
        r.font.name = "Calibri"
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def _slide_number(slide, n, total):
    _add_text(slide, Inches(12.3), Inches(7.05), Inches(1), Inches(0.3),
              f"{n:02d} / {total:02d}", size=10, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)


def _footer_bar(slide):
    from pptx.enum.shapes import MSO_SHAPE
    _add_rect(slide, 0, Inches(7.3), Inches(13.333), Inches(0.2), GOLD)


# ─── Slide builders ────────────────────────────────────────────────────────

def build_cover(prs, n, total, version_label):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(slide, NAVY)

    # Gold accent bar
    _add_rect(slide, 0, Inches(0), Inches(0.3), Inches(7.5), GOLD)

    _add_text(slide, Inches(1.5), Inches(2.5), Inches(11), Inches(1.2),
              "AltaJuris", size=72, color=WHITE, bold=True)
    _add_text(slide, Inches(1.5), Inches(3.5), Inches(11), Inches(0.8),
              "Inteligência jurídica de alto nível", size=28, color=GOLD, bold=False)
    _add_text(slide, Inches(1.5), Inches(4.3), Inches(11), Inches(0.6),
              "Gestão de processos, prazos, documentos e IA — tudo integrado",
              size=18, color=LIGHT_GRAY)

    # Version label
    _add_text(slide, Inches(1.5), Inches(6.5), Inches(11), Inches(0.4),
              version_label, size=12, color=LIGHT_GRAY)
    _slide_number(slide, n, total)
    return slide


def build_bullet_slide(prs, n, total, *, title, subtitle, bullets, accent=GOLD):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(slide, NAVY)

    # Header strip
    _add_rect(slide, 0, 0, Inches(13.333), Inches(0.12), accent)
    _add_text(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
              title, size=34, color=WHITE, bold=True)
    if subtitle:
        _add_text(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.6),
                  subtitle, size=16, color=LIGHT_GRAY)

    _add_bullets(slide, Inches(0.8), Inches(2.3), Inches(11.5), Inches(4.5), bullets,
                 size=18)
    _footer_bar(slide)
    _slide_number(slide, n, total)
    return slide


def build_two_column_slide(prs, n, total, *, title, subtitle, left_label, left_items,
                           right_label, right_items, left_color=EMERALD, right_color=AMBER):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(slide, NAVY)
    _add_rect(slide, 0, 0, Inches(13.333), Inches(0.12), GOLD)

    _add_text(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
              title, size=32, color=WHITE, bold=True)
    if subtitle:
        _add_text(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.6),
                  subtitle, size=15, color=LIGHT_GRAY)

    # Left column
    _add_rect(slide, Inches(0.8), Inches(2.2), Inches(5.8), Inches(4.7),
              RGBColor(0x2A, 0x50, 0x80))
    _add_rect(slide, Inches(0.8), Inches(2.2), Inches(5.8), Inches(0.5), left_color)
    _add_text(slide, Inches(1), Inches(2.25), Inches(5.4), Inches(0.4),
              left_label, size=16, color=WHITE, bold=True)
    _add_bullets(slide, Inches(1), Inches(2.85), Inches(5.4), Inches(3.9), left_items,
                 size=14)

    # Right column
    _add_rect(slide, Inches(6.8), Inches(2.2), Inches(5.8), Inches(4.7),
              RGBColor(0x2A, 0x50, 0x80))
    _add_rect(slide, Inches(6.8), Inches(2.2), Inches(5.8), Inches(0.5), right_color)
    _add_text(slide, Inches(7), Inches(2.25), Inches(5.4), Inches(0.4),
              right_label, size=16, color=WHITE, bold=True)
    _add_bullets(slide, Inches(7), Inches(2.85), Inches(5.4), Inches(3.9), right_items,
                 size=14)

    _footer_bar(slide)
    _slide_number(slide, n, total)
    return slide


def build_section_divider(prs, n, total, *, text, number=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(slide, DARK)
    _add_rect(slide, 0, Inches(3.5), Inches(13.333), Inches(0.04), GOLD)

    if number:
        _add_text(slide, Inches(0.8), Inches(2), Inches(12), Inches(1),
                  number, size=72, color=GOLD, bold=True)
    _add_text(slide, Inches(0.8), Inches(3.6), Inches(12), Inches(1.5),
              text, size=44, color=WHITE, bold=True)
    _slide_number(slide, n, total)
    return slide


def build_closing(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(slide, NAVY)
    _add_rect(slide, 0, Inches(0), Inches(0.3), Inches(7.5), GOLD)

    _add_text(slide, Inches(1.5), Inches(2.3), Inches(11), Inches(1),
              "Pronto para começar?", size=48, color=WHITE, bold=True)
    _add_text(slide, Inches(1.5), Inches(3.4), Inches(11), Inches(0.6),
              "Cadastre seu escritório em altajuris.com.br",
              size=22, color=GOLD)

    _add_text(slide, Inches(1.5), Inches(4.8), Inches(11), Inches(0.4),
              "Site:  https://altajuris.com.br", size=14, color=LIGHT_GRAY)
    _add_text(slide, Inches(1.5), Inches(5.2), Inches(11), Inches(0.4),
              "Dev:   https://desenvolvimento.altajuris.com.br", size=14, color=LIGHT_GRAY)
    _add_text(slide, Inches(1.5), Inches(5.6), Inches(11), Inches(0.4),
              "Docs:  http://development-legaltech-docs-site.s3-website-us-east-1.amazonaws.com",
              size=14, color=LIGHT_GRAY)

    _slide_number(slide, n, total)
    return slide


# ─── Build deck ────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Deck content — mostly hand-curated + a few sections derived from manual
    manual_sections = {
        p.stem: parse_manual(p) for p in MANUAL_DIR.glob("*.md") if p.stem != "index"
    }

    today = _dt.date.today().strftime("%d/%m/%Y")
    version_label = f"Versão atualizada em {today}"

    # Pre-count slides
    slides_plan: list[tuple[str, dict]] = [
        ("cover", {}),
        ("divider", {"number": "01", "text": "O desafio"}),
        ("bullets", {
            "title": "Os desafios do advogado moderno",
            "subtitle": "Por que o AltaJuris existe",
            "bullets": [
                "Consulta manual de andamentos em múltiplos tribunais toma horas por dia",
                "Risco constante de perder prazos processuais",
                "Documentos espalhados em e-mails, pastas locais e WhatsApp",
                "Falta de visão consolidada do escritório: faturamento, casos ativos, produtividade",
                "Peticionamento com certificado digital é complexo e cheio de fricções",
                "Clientes esperam transparência e atualização em tempo real",
            ],
        }),
        ("divider", {"number": "02", "text": "A plataforma"}),
        ("bullets", {
            "title": "AltaJuris — escritório digital completo",
            "subtitle": "Gestão, IA, integrações e conformidade numa plataforma só",
            "bullets": [
                "Gestão de casos com linha de partes (Autor vs Réu) no header",
                "Sincronização automática com eSAJ, PJe e DataJud/CNJ",
                "Controle de prazos com alertas escalonados e countdown visual",
                "Documentos divididos em 'Processo na Íntegra' e 'Documentos e Anexos'",
                "IA para geração de minutas, relatórios e análise de viabilidade",
                "Marketplace conectando cidadãos e advogados especializados",
            ],
        }),
        ("two_column", {
            "title": "Gestão de casos",
            "subtitle": "Cada caso com informações que advogados realmente precisam",
            "left_label": "Cabeçalho do caso",
            "left_items": [
                "Número do processo como título",
                "Linha 'Autor vs Réu' abaixo — identificação rápida",
                "Autor em azul, Réu em vermelho",
                "(+N) com tooltip para múltiplas partes",
                "Badges: status, tribunal, cliente, data",
            ],
            "right_label": "Abas do detalhe",
            "right_items": [
                "Resumo — partes, síntese, cards rápidos",
                "Andamentos — movimentações com badges coloridos",
                "Prazos — lista e calendário com urgência",
                "Timeline — cronologia mensal filtrável",
                "WhatsApp — conversas importadas como evidência",
            ],
        }),
        ("bullets", {
            "title": "Processo na Íntegra vs Documentos e Anexos",
            "subtitle": "Sugestão da Dra. Ana Botan — separação que faz sentido no dia a dia",
            "bullets": [
                "Processo na Íntegra — peças sincronizadas do tribunal (eSAJ/PJe/DataJud)",
                "Ícone do martelo (⚖️) identifica visualmente peças oficiais",
                "Upload bloqueado nesta seção (mantém integridade do autos oficial)",
                "Documentos e Anexos — provas, contratos, áudios, WhatsApp",
                "Classificação automática por IA em 10+ categorias",
                "OCR, assinatura digital e solicitação de assinatura por e-mail",
            ],
        }),
        ("bullets", {
            "title": "Nunca mais perca um prazo",
            "subtitle": "Detecção automática, banner de urgência e notificações escalonadas",
            "bullets": [
                "Detecção automática de prazos em publicações do DJE",
                "Cálculo baseado no CPC (corridos x úteis) por tipo de prazo",
                "Banner vermelho de urgência quando ≤ 1 dia",
                "Calendário visual com cores por status",
                "Notificações escalonadas: 7d, 3d, 1d, no dia",
                "Marcação rápida de conclusão com um clique",
            ],
        }),
        ("bullets", {
            "title": "Importação por OAB em minutos",
            "subtitle": "Trouxemos todos os seus processos de uma só vez",
            "bullets": [
                "Busca paralela em DataJud/CNJ e eSAJ por número da OAB",
                "Suporta múltiplos estados simultaneamente",
                "Polling assíncrono com progresso em tempo real",
                "Dedup por número CNJ — evita duplicatas",
                "Criação automática de casos com partes, vara e comarca",
                "Sincronização periódica para manter andamentos atualizados",
            ],
        }),
        ("divider", {"number": "03", "text": "Inteligência artificial"}),
        ("bullets", {
            "title": "IA trabalhando para você",
            "subtitle": "Claude 3.5 via AWS Bedrock — fundamentada em RAG",
            "bullets": [
                "Chat com os documentos do caso (RAG + embeddings)",
                "Geração de minutas e petições a partir de templates",
                "Relatórios em PDF com análise de viabilidade e estratégia",
                "Extração automática de dados (partes, valores, datas)",
                "Classificação de documentos por tipo em 10+ categorias",
                "Análise de impacto legislativo em casos existentes",
            ],
        }),
        ("bullets", {
            "title": "Construtor de caso (Case Builder)",
            "subtitle": "Wizard guiado para montar processos do zero",
            "bullets": [
                "Etapa 1 — Índice de provas organizado",
                "Etapa 2 — Estimativa de danos calculada por IA",
                "Etapa 3 — Rascunho de petição inicial",
                "Etapa 4 — Exportação em PDF pronta para protocolar",
                "Cross-reference com mensagens do WhatsApp importadas",
            ],
        }),
        ("divider", {"number": "04", "text": "Segurança e conformidade"}),
        ("bullets", {
            "title": "Certificados digitais ICP-Brasil",
            "subtitle": "A1 para automação, A3 para peticionamento",
            "bullets": [
                "Certificado A1 — arquivo .pfx encriptado em repouso (Fernet)",
                "Certificado A3 — token USB via AltaJuris Signer (app próprio gratuito)",
                "Signer suporta SafeSign, SafeNet, Watchdata, OpenSC e outros",
                "Chave privada nunca sai do token — apenas dados públicos",
                "Fallback automático para Lacuna Web PKI quando disponível",
                "Auto-detecção de drivers em Windows e macOS",
            ],
        }),
        ("bullets", {
            "title": "Seus dados protegidos",
            "subtitle": "Arquitetura multi-tenant com isolamento e auditoria",
            "bullets": [
                "Multi-tenant: cada escritório tem dados totalmente isolados",
                "Autenticação via AWS Cognito com política de senha forte",
                "Senhas hashadas (bcrypt) — irrecuperáveis por design",
                "RBAC: admin, advogado, assistente, cliente",
                "Audit log completo de todas as operações sensíveis",
                "Backups automáticos do Aurora PostgreSQL",
            ],
        }),
        ("divider", {"number": "05", "text": "Planos e próximos passos"}),
        ("two_column", {
            "title": "Planos flexíveis",
            "subtitle": "Escolha o plano certo para cada perfil",
            "left_label": "Advogados e escritórios",
            "left_items": [
                "Advogado — R$ 199/mês (1 OAB, até 50 casos)",
                "Escritório — R$ 499/mês (OAB ilimitada, equipe)",
                "Enterprise — sob consulta (SLA, customizações)",
                "Integração com marketplace para leads",
                "Cobrança Asaas (boleto/PIX) para clientes",
            ],
            "right_label": "Cidadãos",
            "right_items": [
                "Básico — R$ 29/mês (1 caso, relatório IA)",
                "Pro — R$ 79/mês (5 casos, análise WhatsApp)",
                "Advanced — R$ 149/mês (ilimitado, marketplace)",
                "Gerador de petições por IA",
                "Score de viabilidade com strengths/weaknesses",
            ],
        }),
        ("bullets", {
            "title": "O que vem por aí",
            "subtitle": "Roadmap próximos meses",
            "bullets": [
                "AltaJuris Signer para macOS e Linux",
                "App mobile (Flutter) com notificações push",
                "Integração com mais tribunais (eProc, TRFs)",
                "Assinatura digital ICP-Brasil integrada (PAdES/CAdES)",
                "Monitoramento legislativo com alertas por embedding",
                "Dashboard executivo com BI (Receita, produtividade, prazos)",
            ],
        }),
        ("closing", {}),
    ]

    total = len(slides_plan)
    for i, (kind, cfg) in enumerate(slides_plan, start=1):
        if kind == "cover":
            build_cover(prs, i, total, version_label)
        elif kind == "divider":
            build_section_divider(prs, i, total, **cfg)
        elif kind == "bullets":
            build_bullet_slide(prs, i, total, **cfg)
        elif kind == "two_column":
            build_two_column_slide(prs, i, total, **cfg)
        elif kind == "closing":
            build_closing(prs, i, total)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    STATIC_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_FILE)
    # Also copy to static/downloads so Docusaurus picks it up during build
    static_file = STATIC_DIR / OUT_FILE.name
    static_file.write_bytes(OUT_FILE.read_bytes())
    print(f"✓ Apresentação gerada: {OUT_FILE}")
    print(f"  Cópia em static: {static_file}")
    print(f"  Slides: {total}")
    print(f"  Data: {today}")
    return OUT_FILE


if __name__ == "__main__":
    build()
